import os
from PIL import Image
from docx import Document
from pptx import Presentation
import pandas as pd
import fitz
import pdfplumber
import openpyxl
import extract_msg
import email
import yaml


def read_rules_file(rules_path, logging):
    """Read rules file and return a DataFrame."""
    try:
        if rules_path.endswith(('.xls', '.xlsx')):
            logging.info(f"Reading Excel file: {rules_path}")
            return pd.read_excel(rules_path)
        elif rules_path.endswith('.csv'):
            logging.info(f"Reading CSV file: {rules_path}")
            return pd.read_csv(rules_path)
        else:
            logging.error(f"Unsupported file format for rules file: {rules_path}")
            return None
    except Exception as e:
        logging.error(f"Error reading rules file: {e}")
        return None


def extract_text_from_pdf(pdf_path, logging):
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text("text")
        if text.strip():
            logging.info("Text extracted using PyMuPDF.")
            return text
    except Exception as e:
        logging.error(f"Error with PyMuPDF: {e}")

    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text()
        if text.strip():
            logging.info("Text extracted using pdfplumber.")
            return text
    except Exception as e:
        logging.error(f"Error with pdfplumber: {e}")

    logging.warning("No text found, attempting OCR...")
    return "No text extracted."


def extract_text_from_ppt(ppt_path, logging):
    presentation = Presentation(ppt_path)
    extracted_text = ""
    for slide_number, slide in enumerate(presentation.slides):
        extracted_text += f"Slide {slide_number + 1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"

        table_text = extract_table_text(slide, logging)
        if table_text:
            extracted_text += "Table Content:\n" + table_text

        extracted_text += "\n"
    logging.info(f"Text extracted from PPT: {ppt_path}")
    return extracted_text


def extract_table_text(slide, logging):
    table_text = ""
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    table_text += cell.text.strip() + " "
                table_text += "\n"
    logging.debug("Table content extracted.")
    return table_text


def extract_text_from_docx(docx_path, logging):
    document = Document(docx_path)
    extracted_text = ""

    for para in document.paragraphs:
        extracted_text += para.text + "\n"

    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                extracted_text += cell.text + " "
            extracted_text += "\n"

    logging.info(f"Text extracted from DOCX: {docx_path}")
    return extracted_text


def list_files_in_directory(directory_path, logging):
    if not os.path.isdir(directory_path):
        logging.error(f"The directory '{directory_path}' does not exist.")
        return []
    logging.info(f"Listing all files in directory: {directory_path}")
    all_files = []
    for root, _, files in os.walk(directory_path):
        for file in files:
            all_files.append(os.path.join(root, file))
    return all_files


def search_words_in_text(text, keywords, logging):
    found_words = {}
    for word in keywords:
        found_words[word] = word in text
    logging.debug(f"Search results for words: {found_words}")
    return found_words

def extract_text_from_excel(file_path):
    if file_path.endswith('.xlsx'):
        return extract_text_from_xlsx(file_path)
    elif file_path.endswith('.xls'):
        return extract_text_from_xls(file_path)
    else:
        raise ValueError("File format not supported. Please provide an .xlsx or .xls file.")

def extract_text_from_xlsx(file_path):
    text = []
    wb = openpyxl.load_workbook(file_path, data_only=True)
    for sheet in wb.sheetnames:
        sheet_text = extract_text_from_sheet(wb[sheet])
        text.extend(sheet_text)
    return text

def extract_text_from_xls(file_path):
    text = []
    df = pd.read_excel(file_path, sheet_name=None)
    for sheet_df in df.values():
        sheet_text = extract_text_from_dataframe(sheet_df)
        text.extend(sheet_text)
    return text

def extract_text_from_sheet(sheet):
    text = []
    for row in sheet.iter_rows():
        for cell in row:
            if isinstance(cell.value, str):
                text.append(cell.value)
    return text

def extract_text_from_dataframe(df):
    text = []
    for value in df.values.flatten():
        if isinstance(value, str):
            text.append(value)
    return text

def extract_text_from_csv(file_path):
    text = []
    df = pd.read_csv(file_path)
    for value in df.values.flatten():
        if isinstance(value, str):
            text.append(value)
    return text

def read_file(file_path, logging):
    logging.info(f"Attempting to read file: {file_path}")
    try:
        if file_path.endswith('.csv'):
            data = pd.read_csv(file_path)
            logging.info("File successfully read as CSV.")
        elif file_path.endswith(('.xls', '.xlsx')):
            data = pd.read_excel(file_path)
            logging.info("File successfully read as Excel.")
        else:
            logging.error("Unsupported file format. Please provide a CSV or Excel file.")
            raise ValueError("Unsupported file format.")
        return data
    except Exception as e:
        logging.error(f"Error reading file: {e}")
        raise


def read_excel_file(file_path, logging):
    logging.info(f"Attempting to read file: {file_path}")
    try:
        if file_path.endswith('.csv'):
            data = pd.read_csv(file_path)
            logging.info("File successfully read as CSV.")
            text = extract_text_from_csv(file_path)
            return text 
        
        elif file_path.endswith(('.xls', '.xlsx')):
            data = None  
            text = extract_text_from_excel(file_path)
            logging.info("File successfully read as Excel.")
            return text  
        
        else:
            logging.error("Unsupported file format. Please provide a CSV or Excel file.")
            raise ValueError("Unsupported file format.")
        
    except Exception as e:
        logging.error(f"Error reading file: {e}")
        raise
    
    
def extract_text_from_log(file_path, logging):
    """
    Extracts all text from the given .log file and logs the process.
    
    Args:
        file_path (str): Path to the .log file.
        logging: Logging object to log the process.
    
    Returns:
        str: The text content from the .log file.
    """
    logging.info(f"Attempting to read log file: {file_path}")
    
    try:
        with open(file_path, 'r') as file:
            content = file.read()  # Read the entire file content
        logging.info("Log file successfully read.")
        return content
    except FileNotFoundError:
        logging.error(f"Error: The file '{file_path}' was not found.")
        return f"Error: The file '{file_path}' was not found."
    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return f"An error occurred: {e}"
    
    
import extract_msg
import logging

def extract_text_from_outlook(file_path, logging):
    """
    Extracts all text from the given .msg (Outlook email) file and logs the process.
    
    Args:
        file_path (str): Path to the .msg (Outlook email) file.
        logging: Logging object to log the process.
    
    Returns:
        str: The text content of the email in the .msg file.
    """
    logging.info(f"Attempting to read Outlook .msg file: {file_path}")
    
    try:
        msg = extract_msg.Message(file_path)
        msg_message = msg.body 
        
        logging.info(f"Outlook .msg file '{file_path}' successfully read.")
        return msg_message
    except FileNotFoundError:
        logging.error(f"Error: The file '{file_path}' was not found.")
        return f"Error: The file '{file_path}' was not found."
    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return f"An error occurred: {e}"

def extract_text_from_eml(file_path, logging):
    """
    Extracts all text from the given .eml (email) file and logs the process.
    
    Args:
        file_path (str): Path to the .eml (email) file.
        logging: Logging object to log the process.
    
    Returns:
        str: The text content of the email in the .eml file.
    """
    logging.info(f"Attempting to read .eml file: {file_path}")
    
    try:
        with open(file_path, 'r') as file:
            msg = email.message_from_file(file)
        
        
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":  
                    email_body = part.get_payload(decode=True).decode(part.get_content_charset())
                    break
        else:
            email_body = msg.get_payload(decode=True).decode(msg.get_content_charset())
        
        logging.info(f".eml file '{file_path}' successfully read.")
        return email_body
    except FileNotFoundError:
        logging.error(f"Error: The file '{file_path}' was not found.")
        return f"Error: The file '{file_path}' was not found."
    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return f"An error occurred: {e}"
    
    
def extract_text_from_yaml(file_path, logging):
    """
    Extracts all text from the given .yaml file and logs the process.
    
    Args:
        file_path (str): Path to the .yaml file.
        logging: Logging object to log the process.
    
    Returns:
        str: The text content of the .yaml file.
    """
    logging.info(f"Attempting to read .yaml file: {file_path}")
    
    try:
        with open(file_path, 'r') as file:
            data = yaml.safe_load(file)
        
        def extract_text(data):
            text = []
            
            if isinstance(data, dict):  
                for key, value in data.items():
                    text.append(str(key))  
                    text.extend(extract_text(value))  
            elif isinstance(data, list): 
                for item in data:
                    text.extend(extract_text(item))  
            else:
                text.append(str(data)) 
            
            return text
        extracted_text = extract_text(data)
        
        logging.info(f".yaml file '{file_path}' successfully read.")
        return "\n".join(extracted_text)
    
    except FileNotFoundError:
        logging.error(f"Error: The file '{file_path}' was not found.")
        return f"Error: The file '{file_path}' was not found."
    except Exception as e:
        logging.error(f"An error occurred: {e}")
        return f"An error occurred: {e}"